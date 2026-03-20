"""Text extraction from PDF, DOCX, PPTX, and ODT files.

Converts binary document formats to plain text for RAG chunking.
Uses stdlib + minimal deps: python-docx for DOCX, pdftotext CLI for PDF.

This module is the "connector" layer — the only part that changes when
moving from local filesystem to S3/data-lake sources.
"""

from __future__ import annotations

import logging
import subprocess
from pathlib import Path
from typing import Optional

log = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".odt", ".txt", ".md", ".xlsx", ".doc"}


def extract_text(path: Path) -> Optional[str]:
    """Extract plain text from a document file.

    Returns None if extraction fails or format is unsupported.
    """
    suffix = path.suffix.lower()

    if suffix in (".md", ".txt"):
        return _read_text_file(path)
    elif suffix == ".pdf":
        return _extract_pdf(path)
    elif suffix == ".docx":
        return _extract_docx(path)
    elif suffix == ".pptx":
        return _extract_pptx(path)
    elif suffix == ".odt":
        return _extract_odt(path)
    elif suffix == ".xlsx":
        return _extract_xlsx(path)
    elif suffix == ".doc":
        return _extract_doc(path)
    else:
        log.debug("Unsupported format: %s", suffix)
        return None


def _read_text_file(path: Path) -> Optional[str]:
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except OSError as e:
        log.warning("Failed to read %s: %s", path, e)
        return None


def _extract_pdf(path: Path) -> Optional[str]:
    """Extract text using pdftotext (poppler-utils) CLI."""
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", str(path), "-"],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        log.warning("pdftotext not installed — install poppler-utils")
    except subprocess.TimeoutExpired:
        log.warning("pdftotext timed out on %s", path)
    except Exception as e:
        log.warning("PDF extraction failed for %s: %s", path, e)

    # Fallback: try pypdf if available
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        text = "\n\n".join(pages)
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        log.warning("pypdf fallback failed for %s: %s", path, e)

    return None


def _extract_docx(path: Path) -> Optional[str]:
    """Extract text from Word documents using python-docx."""
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            return "\n\n".join(paragraphs)
    except ImportError:
        log.warning("python-docx not installed — pip install python-docx")
    except Exception as e:
        log.warning("DOCX extraction failed for %s: %s", path, e)
    return None


def _extract_pptx(path: Path) -> Optional[str]:
    """Extract text from PowerPoint files using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        if slides:
            return "\n\n".join(slides)
    except ImportError:
        log.warning("python-pptx not installed — pip install python-pptx")
    except Exception as e:
        log.warning("PPTX extraction failed for %s: %s", path, e)
    return None


def _extract_odt(path: Path) -> Optional[str]:
    """Extract text from ODT files using zip + XML parsing."""
    import zipfile
    import xml.etree.ElementTree as ET

    try:
        with zipfile.ZipFile(str(path)) as zf:
            with zf.open("content.xml") as f:
                tree = ET.parse(f)
        # Strip all XML tags, keep text
        text = ET.tostring(tree.getroot(), encoding="unicode", method="text")
        if text.strip():
            return text
    except Exception as e:
        log.warning("ODT extraction failed for %s: %s", path, e)
    return None


def _extract_xlsx(path: Path) -> Optional[str]:
    """Extract text from Excel spreadsheets using openpyxl."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(rows))
        wb.close()
        if sheets:
            return "\n\n".join(sheets)
    except ImportError:
        log.warning("openpyxl not installed — pip install openpyxl")
    except Exception as e:
        log.warning("XLSX extraction failed for %s: %s", path, e)
    return None


def _extract_doc(path: Path) -> Optional[str]:
    """Extract text from legacy .doc files via antiword CLI (best-effort)."""
    try:
        result = subprocess.run(
            ["antiword", str(path)],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        log.debug("antiword not installed — .doc files will be skipped")
    except subprocess.TimeoutExpired:
        log.warning("antiword timed out on %s", path)
    except Exception as e:
        log.warning("DOC extraction failed for %s: %s", path, e)
    return None


def extract_directory(
    root: Path,
    output_dir: Optional[Path] = None,
) -> dict[str, str]:
    """Extract text from all supported files in a directory tree.

    Returns dict of {relative_path: extracted_text}.
    If output_dir is set, also writes .txt files there.
    """
    results: dict[str, str] = {}
    supported = 0
    extracted = 0
    failed = 0

    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        if path.name.startswith(".") or "__MACOSX" in str(path):
            continue

        supported += 1
        rel = str(path.relative_to(root))
        text = extract_text(path)

        if text and text.strip():
            results[rel] = text
            extracted += 1

            if output_dir:
                out_path = output_dir / (rel + ".txt")
                out_path.parent.mkdir(parents=True, exist_ok=True)
                out_path.write_text(text, encoding="utf-8")
        else:
            failed += 1
            log.info("No text extracted from %s", rel)

    log.info(
        "Extraction complete: %d supported, %d extracted, %d failed",
        supported, extracted, failed,
    )
    return results
